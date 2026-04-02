"""
PPTX转图片脚本
将PPTX文件的每张幻灯片转换为PNG图片
使用python-pptx提取图片+文字信息，再用Pillow渲染为图片
"""
import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont
import io
import zipfile

def extract_slide_images_from_pptx(pptx_path, output_dir):
    """从PPTX文件提取幻灯片，转换为图片"""
    os.makedirs(output_dir, exist_ok=True)
    
    prs = Presentation(pptx_path)
    slide_count = len(prs.slides)
    slide_files = []
    
    print(f"PPTX共有 {slide_count} 张幻灯片")
    
    # 方法：直接从PPTX ZIP包中提取嵌入的图片，并生成幻灯片预览
    # PPTX本质是ZIP包，里面有media文件夹存放图片
    
    with zipfile.ZipFile(pptx_path, 'r') as z:
        media_files = [f for f in z.namelist() if f.startswith('ppt/media/')]
        print(f"PPTX包含 {len(media_files)} 个媒体文件")
    
    # 逐张幻灯片处理
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        
        # 创建幻灯片画布 (16:9 = 1280x720)
        slide_w = int(prs.slide_width.inches * 96)
        slide_h = int(prs.slide_height.inches * 96)
        # 统一缩放到1280x720
        canvas_w, canvas_h = 1280, 720
        scale = canvas_w / slide_w
        
        img = Image.new('RGB', (canvas_w, canvas_h), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)
        
        # 提取图片形状并绘制
        has_image = False
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    pic_data = shape.image.blob
                    pic_img = Image.open(io.BytesIO(pic_data))
                    # 计算位置和尺寸
                    x = int(shape.left * scale / 914400 * 96)
                    y = int(shape.top * scale / 914400 * 96)
                    w = int(shape.width * scale / 914400 * 96)
                    h = int(shape.height * scale / 914400 * 96)
                    # 确保不超出画布
                    x = max(0, min(x, canvas_w - 1))
                    y = max(0, min(y, canvas_h - 1))
                    w = max(1, min(w, canvas_w - x))
                    h = max(1, min(h, canvas_h - y))
                    pic_img_resized = pic_img.resize((w, h), Image.LANCZOS)
                    img.paste(pic_img_resized, (x, y))
                    has_image = True
                except Exception as e:
                    print(f"  图片形状处理失败: {e}")
        
        # 绘制文本（简单渲染）
        try:
            font_title = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 28)
            font_body = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 16)
        except:
            font_title = ImageFont.load_default()
            font_body = ImageFont.load_default()
        
        text_y = 30
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                # 跳过图片形状的文字
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    continue
                text = shape.text.strip()
                if not text:
                    continue
                
                x_pos = int(shape.left * scale / 914400 * 96) if shape.left else 50
                y_pos = int(shape.top * scale / 914400 * 96) if shape.top else text_y
                x_pos = max(20, min(x_pos, canvas_w - 200))
                y_pos = max(20, min(y_pos, canvas_h - 40))
                
                # 判断是否是标题（位置靠上，或字体较大）
                is_title = y_pos < 150 or (hasattr(shape, 'placeholder_format') and 
                           shape.placeholder_format and 
                           shape.placeholder_format.idx == 0)
                
                font = font_title if is_title else font_body
                color = (51, 51, 51)
                
                # 分行绘制（每行最多40个字符）
                lines = text.split('\n')
                for line in lines[:8]:  # 最多8行
                    if line.strip():
                        # 截断过长的行
                        if len(line) > 50:
                            line = line[:50] + '...'
                        try:
                            draw.text((x_pos, y_pos), line, font=font, fill=color)
                        except Exception:
                            pass
                        y_pos += (font_title.size if is_title else font_body.size) + 4
        
        # 如果幻灯片完全是图片（背景），直接用第一张图片作为幻灯片
        if not has_image:
            # 添加幻灯片编号标记
            draw.rectangle([0, canvas_h-30, canvas_w, canvas_h], fill=(240,240,240))
            draw.text((10, canvas_h-25), f"第 {slide_num} 张", font=font_body, fill=(120,120,120))
        
        # 保存
        out_path = os.path.join(output_dir, f'slide_{slide_num:03d}.jpg')
        img.save(out_path, 'JPEG', quality=85)
        slide_files.append(f'slide_{slide_num:03d}.jpg')
        print(f"  已处理第 {slide_num} 张 -> {out_path}")
    
    return slide_files


if __name__ == '__main__':
    pptx_path = r"C:\Users\BJ7070\WorkBuddy\20260324105325\industry-knowledge-base\files\minutes\2026年4月2日 审计板块行业线经营例会纪要.pptx"
    output_dir = r"C:\Users\BJ7070\WorkBuddy\20260324105325\industry-knowledge-base\files\pptx-slides\20260402-meeting"
    
    slides = extract_slide_images_from_pptx(pptx_path, output_dir)
    
    # 保存幻灯片列表JSON
    meta = {
        "title": "2026年4月2日 审计板块行业线经营例会纪要",
        "date": "2026-04-02",
        "total": len(slides),
        "slides": slides
    }
    with open(os.path.join(output_dir, 'meta.json'), 'w', encoding='utf-8') as f:
        json.dump(meta, f, ensure_ascii=False, indent=2)
    
    print(f"\n完成！共 {len(slides)} 张幻灯片，已保存到 {output_dir}")
